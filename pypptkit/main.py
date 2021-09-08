"""
main entry point
"""
import pylogconf.core
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pytconf import register_endpoint, get_free_args, register_main, config_arg_parse_and_launch
from pyvardump import dump_print

from pypptkit.static import DESCRIPTION, APP_NAME, VERSION_STR
from pypptkit.utils import get_sorted_refs, download


@register_endpoint(
    allow_free_args=True,
    description="Extract text from ppt files",
)
def extract_text() -> None:
    for filename in get_free_args():
        presentation = Presentation(filename)
        for slide_number, slide in enumerate(presentation.slides):
            print(f"slide number {slide_number}")
            if slide.name != "":
                print(f"slide.name {slide.name}")
            # This block extracts hyperlinks
            for v in slide.part.rels.values():
                target: str = v.target_ref
                if target.startswith(".."):
                    continue
                print(v.target_ref)
            for shape in slide.shapes:
                # pylint: disable=no-member
                if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    print(shape.shape_type)
                    print(shape.text)
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    print(shape.shape_type)
                    print(shape.text)
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    print(shape.shape_type)
                    for row in shape.table.rows:
                        for cell in row.cells:
                            print(cell.text)


@register_endpoint(
    allow_free_args=True,
    description="Extract links from ppt files",
)
def extract_links() -> None:
    refs = get_sorted_refs(get_free_args())
    for ref in refs:
        print(ref)


@register_endpoint(
    allow_free_args=True,
    description="Download links from ppt files",
)
def download_links() -> None:
    refs = get_sorted_refs(get_free_args())
    for ref in refs:
        download(ref)


@register_endpoint(
    allow_free_args=True,
    description="Dump object of ppt files",
)
def dump_slides() -> None:
    for filename in get_free_args():
        print(f"filename {filename}")
        presentation = Presentation(filename)
        for number, slide in enumerate(presentation.slides):
            print(f"slide {number}")
            dump_print(slide)


@register_main(
    main_description=DESCRIPTION,
    app_name=APP_NAME,
    version=VERSION_STR,
)
def main():
    pylogconf.core.setup()
    config_arg_parse_and_launch()


if __name__ == '__main__':
    main()
