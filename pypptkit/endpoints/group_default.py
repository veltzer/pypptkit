"""
The default group of operations that pypptkit has
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pytconf import register_endpoint, register_function_group, get_free_args
from pyvardump import dump

import pypptkit
import pypptkit.version

GROUP_NAME_DEFAULT = "default"
GROUP_DESCRIPTION_DEFAULT = "all pypptkit commands"


def register_group_default():
    """
    register the name and description of this group
    """
    register_function_group(
        function_group_name=GROUP_NAME_DEFAULT,
        function_group_description=GROUP_DESCRIPTION_DEFAULT,
    )


@register_endpoint(
    group=GROUP_NAME_DEFAULT,
)
def version() -> None:
    """
    Print version
    """
    print(pypptkit.version.VERSION_STR)


@register_endpoint(
    group=GROUP_NAME_DEFAULT,
    allow_free_args=True,
)
def extract_text() -> None:
    """
    extract text from ppt files
    """
    for filename in get_free_args():
        presentation = Presentation(filename)
        for slide_number, slide in enumerate(presentation.slides):
            print(f"slide number {slide_number}")
            if slide.name != "":
                print(f"slide.name {slide.name}")
            # This block extracts hyperlinks
            for v in slide.part.rels.values():
                target: str = v.target_ref
                if target.startswith(".."):
                    continue
                print(v.target_ref)
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    print(shape.shape_type)
                    print(shape.text)
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    print(shape.shape_type)
                    print(shape.text)
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    print(shape.shape_type)
                    for row in shape.table.rows:
                        for cell in row.cells:
                            print(cell.text)


@register_endpoint(
    group=GROUP_NAME_DEFAULT,
    allow_free_args=True,
)
def extract_links() -> None:
    """
    extract links from ppt files
    """
    refs = set()
    for filename in get_free_args():
        presentation = Presentation(filename)
        for slide in presentation.slides:
            for v in slide.part.rels.values():
                target: str = v.target_ref
                if target.startswith(".."):
                    continue
                refs.add(target)
    refs = sorted(list(refs))
    for ref in refs:
        print(ref)


@register_endpoint(
    group=GROUP_NAME_DEFAULT,
    allow_free_args=True,
)
def download_links() -> None:
    """
    download links from ppt files
    """
    refs = set()
    for filename in get_free_args():
        presentation = Presentation(filename)
        for slide in presentation.slides:
            for v in slide.part.rels.values():
                target: str = v.target_ref
                if target.startswith(".."):
                    continue
                refs.add(target)
    refs = sorted(list(refs))
    for ref in refs:
        print(ref)


@register_endpoint(
    group=GROUP_NAME_DEFAULT,
    allow_free_args=True,
)
def dump_slides() -> None:
    """
    dump object of ppt files
    """
    for filename in get_free_args():
        print(f"filename {filename}")
        presentation = Presentation(filename)
        for number, slide in enumerate(presentation.slides):
            print(f"slide {number}")
            dump(slide)
